"""
ë©”ì¸ ì‹¤í–‰ ëª¨ë“ˆ (Main Entry Point)

í•œí™”íì…€ ë²ˆì—­ í”„ë¡œì íŠ¸ì˜ ë©”ì¸ ì‹¤í–‰ íŒŒì¼ì…ë‹ˆë‹¤.
Google Sheetsì—ì„œ ì‘ì—…ì„ ê°€ì ¸ì™€ ìˆœì°¨ì ìœ¼ë¡œ ì²˜ë¦¬í•©ë‹ˆë‹¤.
"""

import os
import shutil
import traceback

from docx import Document
from pptx import Presentation

from .config import (
    ORIGIN_FOLDER,
    COMPLETED_FOLDER,
    SUPPORTED_EXTENSIONS,
    ALL_SUPPORTED_EXTENSIONS,
    validate_config
)
from .translator import generate_context
from .handlers import process_docx, process_pptx, process_xlsx
from .sheets_manager import SheetsManager, Status
from .converter import convert_doc_to_docx, needs_conversion, get_converted_extension
from .slack_notifier import send_completion_notification, send_error_notification
from .glossary import get_glossary


def normalize_extension(file_name):
    """
    íŒŒì¼ëª…ì˜ í™•ì¥ìë¥¼ ì†Œë¬¸ìë¡œ ë³€í™˜í•©ë‹ˆë‹¤.
    
    Args:
        file_name (str): íŒŒì¼ëª… (í™•ì¥ì í¬í•¨)
        
    Returns:
        tuple: (ì •ê·œí™”ëœ íŒŒì¼ëª…, ì›ë³¸ í™•ì¥ìê°€ ëŒ€ë¬¸ìì˜€ëŠ”ì§€ ì—¬ë¶€)
    """
    name, ext = os.path.splitext(file_name)
    ext_lower = ext.lower()
    
    # í™•ì¥ìê°€ ëŒ€ë¬¸ìì˜€ëŠ”ì§€ í™•ì¸
    was_uppercase = (ext != ext_lower)
    
    # ì†Œë¬¸ìë¡œ ë³€í™˜ëœ íŒŒì¼ëª…
    normalized_name = f"{name}{ext_lower}"
    
    return normalized_name, was_uppercase


def build_file_path(upper_path, sub_path, file_name):
    """
    ìƒìœ„ê²½ë¡œ, ì„¸ë¶€ê²½ë¡œ, íŒŒì¼ëª…ì„ ì¡°í•©í•˜ì—¬ ì „ì²´ íŒŒì¼ ê²½ë¡œë¥¼ ìƒì„±í•©ë‹ˆë‹¤.
    
    .doc íŒŒì¼ì˜ ê²½ìš° ì‘ì—… íŒŒì¼ì€ .docxë¡œ ìƒì„±ë©ë‹ˆë‹¤.
    í™•ì¥ìê°€ ëŒ€ë¬¸ìì¸ ê²½ìš° ì†Œë¬¸ìë¡œ ë³€í™˜í•©ë‹ˆë‹¤.
    
    Args:
        upper_path (str): ìƒìœ„ ê²½ë¡œ (ì˜ˆ: "MC")
        sub_path (str): ì„¸ë¶€ ê²½ë¡œ (ì˜ˆ: "10.ë¶„ì„ë‹¨ê³„")
        file_name (str): íŒŒì¼ëª… (í™•ì¥ì í¬í•¨)
        
    Returns:
        tuple: (ì›ë³¸ íŒŒì¼ ê²½ë¡œ, ì™„ë£Œ í´ë” ê²½ë¡œ, ì›ë³¸ ë³µì‚¬ë³¸ ê²½ë¡œ, ì‘ì—… íŒŒì¼ ê²½ë¡œ, ì •ê·œí™”ëœ íŒŒì¼ëª…)
    """
    # ì›ë³¸ íŒŒì¼ ê²½ë¡œ (ì›ë³¸ íŒŒì¼ëª… ê·¸ëŒ€ë¡œ ì‚¬ìš©)
    origin_path = os.path.join(ORIGIN_FOLDER, upper_path, sub_path, file_name)
    
    # ì™„ë£Œ í´ë” ë‚´ ê²½ë¡œ (ë™ì¼ êµ¬ì¡° ìœ ì§€)
    completed_dir = os.path.join(COMPLETED_FOLDER, upper_path, sub_path)
    
    # í™•ì¥ì ì†Œë¬¸ìë¡œ ì •ê·œí™”
    normalized_file_name, _ = normalize_extension(file_name)
    
    # ì›ë³¸ ë³µì‚¬ë³¸ ê²½ë¡œ (í™•ì¥ì ì†Œë¬¸ìë¡œ ì €ì¥)
    completed_original = os.path.join(completed_dir, normalized_file_name)
    
    # ì‘ì—… íŒŒì¼ ê²½ë¡œ ê²°ì •
    name, ext = os.path.splitext(normalized_file_name)
    
    # .doc íŒŒì¼ì€ .docxë¡œ ë³€í™˜í•˜ì—¬ ì‘ì—…
    if ext.lower() == '.doc':
        work_file_name = f"{name} - en.docx"  # .doc â†’ .docx ë³€í™˜
    else:
        work_file_name = f"{name} - en{ext}"
    
    work_file_path = os.path.join(completed_dir, work_file_name)
    
    return origin_path, completed_dir, completed_original, work_file_path, normalized_file_name


def prepare_work_files(origin_path, completed_dir, completed_original, work_file_path):
    """
    ì‘ì—… íŒŒì¼ì„ ì¤€ë¹„í•©ë‹ˆë‹¤.
    - ì™„ë£Œ í´ë”ì— ë™ì¼ ê²½ë¡œ ìƒì„±
    - ì›ë³¸ íŒŒì¼ ë³µì‚¬ (ë°±ì—…)
    - ì‘ì—… íŒŒì¼ ìƒì„± (ë²ˆì—­ ëŒ€ìƒ)
    - .doc íŒŒì¼ì€ .docxë¡œ ë³€í™˜
    
    Args:
        origin_path (str): ì›ë³¸ íŒŒì¼ ê²½ë¡œ
        completed_dir (str): ì™„ë£Œ í´ë” ê²½ë¡œ
        completed_original (str): ì›ë³¸ ë³µì‚¬ë³¸ ê²½ë¡œ
        work_file_path (str): ì‘ì—… íŒŒì¼ ê²½ë¡œ
        
    Returns:
        str: ì‹¤ì œ ì‘ì—… íŒŒì¼ ê²½ë¡œ (ì„±ê³µ ì‹œ)
        None: ì‹¤íŒ¨ ì‹œ
    """
    try:
        # 1. ì™„ë£Œ í´ë”ì— ë™ì¼ ê²½ë¡œ ìƒì„±
        os.makedirs(completed_dir, exist_ok=True)
        
        # 2. ì›ë³¸ íŒŒì¼ ë³µì‚¬ (ë°±ì—…ìš© - ì›ë³¸ í˜•ì‹ ê·¸ëŒ€ë¡œ)
        if not os.path.exists(completed_original):
            shutil.copy2(origin_path, completed_original)
            print(f"   ğŸ“ ì›ë³¸ ë³µì‚¬ ì™„ë£Œ: {os.path.basename(completed_original)}")
        
        # 3. íŒŒì¼ í˜•ì‹ì— ë”°ë¥¸ ì‘ì—… íŒŒì¼ ìƒì„±
        ext = os.path.splitext(origin_path)[1].lower()
        
        if ext == '.doc':
            # .doc â†’ .docx ë³€í™˜
            # ë¨¼ì € ì›ë³¸ì„ ì„ì‹œë¡œ ë³µì‚¬í•œ í›„ ë³€í™˜
            temp_doc_path = os.path.join(completed_dir, os.path.basename(origin_path))
            if not os.path.exists(temp_doc_path):
                shutil.copy2(origin_path, temp_doc_path)
            
            # .docxë¡œ ë³€í™˜
            convert_doc_to_docx(temp_doc_path, work_file_path)
            
            # ì„ì‹œ .doc íŒŒì¼ ì‚­ì œ (ì›ë³¸ ë³µì‚¬ë³¸ì´ ì´ë¯¸ ìˆìœ¼ë¯€ë¡œ)
            if temp_doc_path != completed_original:
                try:
                    os.remove(temp_doc_path)
                except:
                    pass
                    
            print(f"   ğŸ“ ì‘ì—… íŒŒì¼ ìƒì„± (ë³€í™˜ë¨): {os.path.basename(work_file_path)}")
        else:
            # ë‹¤ë¥¸ í˜•ì‹ì€ ê·¸ëŒ€ë¡œ ë³µì‚¬
            shutil.copy2(origin_path, work_file_path)
            print(f"   ğŸ“ ì‘ì—… íŒŒì¼ ìƒì„±: {os.path.basename(work_file_path)}")
        
        return work_file_path
        
    except Exception as e:
        print(f"   âŒ íŒŒì¼ ì¤€ë¹„ ì‹¤íŒ¨: {e}")
        return None


def extract_sample_text(file_path):
    """
    íŒŒì¼ì—ì„œ Context ë¶„ì„ìš© ìƒ˜í”Œ í…ìŠ¤íŠ¸ë¥¼ ì¶”ì¶œí•©ë‹ˆë‹¤.
    
    Args:
        file_path (str): íŒŒì¼ ê²½ë¡œ
        
    Returns:
        str: ì¶”ì¶œëœ ìƒ˜í”Œ í…ìŠ¤íŠ¸
    """
    sample_text = ""
    
    # í™•ì¥ìë¥¼ ì†Œë¬¸ìë¡œ ë³€í™˜í•˜ì—¬ ëŒ€ì†Œë¬¸ì êµ¬ë¶„ ì—†ì´ ì²˜ë¦¬
    file_path_lower = file_path.lower()
    
    if file_path_lower.endswith('.docx'):
        doc = Document(file_path)
        sample_text = "\n".join([p.text for p in doc.paragraphs[:300]])
        
    elif file_path_lower.endswith('.pptx'):
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            if i >= 3:
                break
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    sample_text += shape.text + "\n"
                    
    elif file_path_lower.endswith('.xlsx'):
        sample_text = "MES Excel Data"
    
    return sample_text


def process_single_file(work_file_path, file_context, sheets_manager=None, row_index=None):
    """
    ë‹¨ì¼ íŒŒì¼ì„ ë²ˆì—­ ì²˜ë¦¬í•©ë‹ˆë‹¤.
    
    Args:
        work_file_path (str): ì‘ì—… íŒŒì¼ ê²½ë¡œ
        file_context (str): ë²ˆì—­ ì§€ì¹¨ (Context)
        sheets_manager (SheetsManager, optional): ì‹œíŠ¸ ê´€ë¦¬ì (ì§„í–‰ ìƒí™© ì¶”ì ìš©)
        row_index (int, optional): ì‹œíŠ¸ í–‰ ë²ˆí˜¸
        
    Returns:
        str: ë²ˆì—­ëœ íŒŒì¼ ê²½ë¡œ (ì„±ê³µ ì‹œ)
        None: ì‹¤íŒ¨ ì‹œ
    """
    # í™•ì¥ìë¥¼ ì†Œë¬¸ìë¡œ ë³€í™˜í•˜ì—¬ ëŒ€ì†Œë¬¸ì êµ¬ë¶„ ì—†ì´ ì²˜ë¦¬
    file_path_lower = work_file_path.lower()
    
    if file_path_lower.endswith('.docx'):
        return process_docx(work_file_path, file_context, sheets_manager, row_index)
    elif file_path_lower.endswith('.pptx'):
        return process_pptx(work_file_path, file_context, sheets_manager, row_index)
    elif file_path_lower.endswith('.xlsx'):
        return process_xlsx(work_file_path, file_context, sheets_manager, row_index)
    else:
        print(f"   âš ï¸ ì§€ì›í•˜ì§€ ì•ŠëŠ” íŒŒì¼ í˜•ì‹: {work_file_path}")
        return None


def process_task(sheets_manager, task):
    """
    ë‹¨ì¼ ì‘ì—…ì„ ì²˜ë¦¬í•©ë‹ˆë‹¤.
    
    - ìƒíƒœê°€ "ëŒ€ê¸°"ì¸ ê²½ìš°: íŒŒì¼ ë³µì‚¬ í›„ ë²ˆì—­ ì‹œì‘
    - ìƒíƒœê°€ "ì§„í–‰ì¤‘" ë˜ëŠ” "ì˜¤ë¥˜"ì¸ ê²½ìš°: ê¸°ì¡´ "-en" íŒŒì¼ë¡œ ì´ì–´ì„œ ë²ˆì—­
    
    Args:
        sheets_manager (SheetsManager): ì‹œíŠ¸ ê´€ë¦¬ì
        task (dict): ì‘ì—… ì •ë³´
        
    Returns:
        bool: ì„±ê³µ ì—¬ë¶€
    """
    row_index = task['row_index']
    upper_path = task['upper_path']
    sub_path = task['sub_path']
    file_name = task['file_name']
    current_status = task.get('status', 'ëŒ€ê¸°')  # í˜„ì¬ ìƒíƒœ
    
    # ì´ì–´í•˜ê¸° ëª¨ë“œ ì—¬ë¶€
    is_resume_mode = current_status in ['ì§„í–‰ì¤‘', 'ì˜¤ë¥˜']
    
    print(f"\n{'='*60}")
    print(f"ğŸ“„ íŒŒì¼ ì²˜ë¦¬ ì‹œì‘: {file_name}")
    print(f"   ê²½ë¡œ: {upper_path}/{sub_path}")
    if is_resume_mode:
        print(f"   ğŸ”„ ì´ì–´í•˜ê¸° ëª¨ë“œ (ì´ì „ ìƒíƒœ: {current_status})")
    print(f"{'='*60}")
    
    try:
        # 1. íŒŒì¼ ê²½ë¡œ êµ¬ì„± (í™•ì¥ì ì†Œë¬¸ìë¡œ ì •ê·œí™”)
        origin_path, completed_dir, completed_original, work_file_path, normalized_file_name = build_file_path(
            upper_path, sub_path, file_name
        )
        
        # 2. í™•ì¥ìê°€ ëŒ€ë¬¸ìì˜€ìœ¼ë©´ Google Sheetsì—ì„œ íŒŒì¼ëª… ì—…ë°ì´íŠ¸
        if file_name != normalized_file_name:
            print(f"   ğŸ“ íŒŒì¼ëª… í™•ì¥ì ì •ê·œí™”: {file_name} â†’ {normalized_file_name}")
            sheets_manager.update_file_name(row_index, normalized_file_name)
            file_name = normalized_file_name  # ì´í›„ ë¡œì§ì—ì„œ ì‚¬ìš©í•  íŒŒì¼ëª… ì—…ë°ì´íŠ¸
        
        # 3. ì›ë³¸ íŒŒì¼ ì¡´ì¬ í™•ì¸
        if not os.path.exists(origin_path):
            raise FileNotFoundError(f"ì›ë³¸ íŒŒì¼ì„ ì°¾ì„ ìˆ˜ ì—†ìŠµë‹ˆë‹¤: {origin_path}")
        
        # 4. íŒŒì¼ í™•ì¥ì í™•ì¸ (ìµœì‹  í˜•ì‹ + ë³€í™˜ ê°€ëŠ¥ í˜•ì‹ ëª¨ë‘ í—ˆìš©)
        ext = os.path.splitext(file_name)[1].lower()
        if ext not in ALL_SUPPORTED_EXTENSIONS:
            raise ValueError(f"ì§€ì›í•˜ì§€ ì•ŠëŠ” íŒŒì¼ í˜•ì‹ì…ë‹ˆë‹¤: {ext}")
        
        # 5. ì›ë³¸ íŒŒì¼ ë¬´ê²°ì„± ê²€ì‚¬ (ì†ìƒëœ íŒŒì¼ ì¡°ê¸° ê°ì§€)
        if not verify_file_integrity(origin_path):
            raise ValueError(f"ì›ë³¸ íŒŒì¼ì´ ì†ìƒë˜ì—ˆìŠµë‹ˆë‹¤ (ì—´ ìˆ˜ ì—†ìŒ): {origin_path}")
        
        # 5. ì§„í–‰ìƒíƒœ 'ì§„í–‰ì¤‘'ìœ¼ë¡œ ë³€ê²½ + ì‹œì‘ì‹œê°„ ê¸°ë¡
        # ì´ì–´í•˜ê¸° ëª¨ë“œì—ì„œëŠ” í† í° ì´ˆê¸°í™” í•˜ì§€ ì•ŠìŒ
        if is_resume_mode:
            sheets_manager.update_status(row_index, 'ì§„í–‰ì¤‘')
            print(f"   âœ… ìƒíƒœ ë³€ê²½: ì§„í–‰ì¤‘ (ì´ì–´í•˜ê¸°)")
        else:
            sheets_manager.start_task(row_index)
            print(f"   âœ… ìƒíƒœ ë³€ê²½: ì§„í–‰ì¤‘")
        
        # 6. ì‘ì—… íŒŒì¼ ì¤€ë¹„
        if is_resume_mode:
            # ì´ì–´í•˜ê¸° ëª¨ë“œ: ê¸°ì¡´ "-en" íŒŒì¼ ì‚¬ìš©
            actual_work_path = prepare_work_files_resume(work_file_path, origin_path, completed_dir, completed_original)
        else:
            # ìƒˆ ì‘ì—…: íŒŒì¼ ë³µì‚¬ í›„ ì‹œì‘
            actual_work_path = prepare_work_files(origin_path, completed_dir, completed_original, work_file_path)
        
        if not actual_work_path:
            raise Exception("ì‘ì—… íŒŒì¼ ì¤€ë¹„ ì‹¤íŒ¨")
        
        # 7. Context ë¶„ì„
        print(f"   ğŸ¤– Context ë¶„ì„ ì¤‘...")
        sample_text = extract_sample_text(actual_work_path)
        file_context = generate_context(sample_text)
        print(f"   âœ… Context ë¶„ì„ ì™„ë£Œ")
        
        # 8. ë²ˆì—­ ì‹¤í–‰ (ì‹œíŠ¸ ì§„í–‰ ìƒí™© ì¶”ì  í¬í•¨)
        result = process_single_file(actual_work_path, file_context, sheets_manager, row_index)
        
        if result:
            # 9. ì™„ë£Œ ì²˜ë¦¬
            sheets_manager.mark_completed(row_index)
            print(f"\n   ğŸ‰ ë²ˆì—­ ì™„ë£Œ!")
            
            # 10. Slack ì™„ë£Œ ì•Œë¦¼ ì „ì†¡
            try:
                times = sheets_manager.get_task_times(row_index)
                progress = sheets_manager.get_overall_progress()
                file_path = f"{upper_path}/{sub_path}"
                
                send_completion_notification(
                    file_name=file_name,
                    file_path=file_path,
                    start_time=times['start_time'],
                    end_time=times['end_time'],
                    progress_percent=progress
                )
            except Exception as slack_error:
                print(f"   âš ï¸ Slack ì•Œë¦¼ ì „ì†¡ ì‹¤íŒ¨: {slack_error}")
            
            return True
        else:
            raise Exception("ë²ˆì—­ ì²˜ë¦¬ ì‹¤íŒ¨")
            
    except Exception as e:
        # ì˜¤ë¥˜ ë°œìƒ ì‹œ ê¸°ë¡
        error_msg = str(e)
        error_trace = traceback.format_exc()
        module_name = "main.process_task"
        
        # ìƒì„¸ ì˜¤ë¥˜ ë©”ì‹œì§€ êµ¬ì„±
        detailed_error = f"{error_msg}\n\nìƒì„¸:\n{error_trace}"
        
        print(f"\n   âŒ ì˜¤ë¥˜ ë°œìƒ: {error_msg}")
        sheets_manager.record_error(row_index, detailed_error, module_name)
        
        # Slack ì˜¤ë¥˜ ì•Œë¦¼ ì „ì†¡
        try:
            slack_error_msg = f"*íŒŒì¼*: {file_name}\n*ê²½ë¡œ*: {upper_path}/{sub_path}\n*ì˜¤ë¥˜*: {error_msg}"
            send_error_notification(slack_error_msg)
        except Exception as slack_error:
            print(f"   âš ï¸ Slack ì•Œë¦¼ ì „ì†¡ ì‹¤íŒ¨: {slack_error}")
        
        return False


def verify_file_integrity(file_path):
    """
    íŒŒì¼ì˜ ë¬´ê²°ì„±ì„ ê²€ì‚¬í•©ë‹ˆë‹¤.
    
    ì†ìƒëœ íŒŒì¼(ë¹ˆ íŒŒì¼, ì—´ë¦¬ì§€ ì•ŠëŠ” íŒŒì¼)ì¸ì§€ í™•ì¸í•©ë‹ˆë‹¤.
    
    Args:
        file_path (str): ê²€ì‚¬í•  íŒŒì¼ ê²½ë¡œ
        
    Returns:
        bool: íŒŒì¼ì´ ì •ìƒì´ë©´ True, ì†ìƒë˜ì—ˆìœ¼ë©´ False
    """
    try:
        # íŒŒì¼ í¬ê¸° í™•ì¸ (0ë°”ì´íŠ¸ë©´ ì†ìƒ)
        if os.path.getsize(file_path) == 0:
            return False
        
        # í™•ì¥ìë³„ ë¬´ê²°ì„± ê²€ì‚¬
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.docx':
            # python-docxë¡œ ì—´ì–´ë³´ê¸°
            doc = Document(file_path)
            # ìµœì†Œí•œ bodyê°€ ìˆëŠ”ì§€ í™•ì¸
            _ = doc.paragraphs
            return True
            
        elif ext == '.pptx':
            # python-pptxë¡œ ì—´ì–´ë³´ê¸°
            prs = Presentation(file_path)
            _ = prs.slides
            return True
            
        elif ext == '.xlsx':
            # íŒŒì¼ í¬ê¸°ë§Œ í™•ì¸ (xlwingsëŠ” Excel í•„ìš”)
            return os.path.getsize(file_path) > 1000  # ìµœì†Œ 1KB
            
        return True  # ë‹¤ë¥¸ í˜•ì‹ì€ ê¸°ë³¸ì ìœ¼ë¡œ í†µê³¼
        
    except Exception as e:
        print(f"   âš ï¸ íŒŒì¼ ë¬´ê²°ì„± ê²€ì‚¬ ì‹¤íŒ¨: {e}")
        return False


def prepare_work_files_resume(work_file_path, origin_path, completed_dir, completed_original):
    """
    ì´ì–´í•˜ê¸° ëª¨ë“œì—ì„œ ì‘ì—… íŒŒì¼ì„ ì¤€ë¹„í•©ë‹ˆë‹¤.
    
    ê¸°ì¡´ "-en" íŒŒì¼ì´ ìˆìœ¼ë©´ ë¬´ê²°ì„± ê²€ì‚¬ í›„ ì‚¬ìš©í•˜ê³ ,
    ì†ìƒë˜ì—ˆê±°ë‚˜ ì—†ìœ¼ë©´ ìƒˆë¡œ ìƒì„±í•©ë‹ˆë‹¤.
    
    Args:
        work_file_path (str): ì‘ì—… íŒŒì¼ ê²½ë¡œ ("-en" íŒŒì¼)
        origin_path (str): ì›ë³¸ íŒŒì¼ ê²½ë¡œ
        completed_dir (str): ì™„ë£Œ í´ë” ê²½ë¡œ
        completed_original (str): ì›ë³¸ ë³µì‚¬ë³¸ ê²½ë¡œ
        
    Returns:
        str: ì‹¤ì œ ì‘ì—… íŒŒì¼ ê²½ë¡œ (ì„±ê³µ ì‹œ)
        None: ì‹¤íŒ¨ ì‹œ
    """
    try:
        # 1. "-en" íŒŒì¼ì´ ì´ë¯¸ ì¡´ì¬í•˜ëŠ”ì§€ í™•ì¸
        if os.path.exists(work_file_path):
            print(f"   âœ… ê¸°ì¡´ ì‘ì—… íŒŒì¼ ë°œê²¬: {os.path.basename(work_file_path)}")
            
            # 2. íŒŒì¼ ë¬´ê²°ì„± ê²€ì‚¬
            if verify_file_integrity(work_file_path):
                print(f"   ğŸ”„ ì´ì–´ì„œ ë²ˆì—­ì„ ì§„í–‰í•©ë‹ˆë‹¤...")
                return work_file_path
            else:
                # ì†ìƒëœ íŒŒì¼ ì‚­ì œ í›„ ìƒˆë¡œ ìƒì„±
                print(f"   âš ï¸ ê¸°ì¡´ íŒŒì¼ì´ ì†ìƒë˜ì—ˆìŠµë‹ˆë‹¤. ì‚­ì œ í›„ ìƒˆë¡œ ìƒì„±í•©ë‹ˆë‹¤...")
                try:
                    os.remove(work_file_path)
                except Exception as del_err:
                    print(f"   âŒ ì†ìƒëœ íŒŒì¼ ì‚­ì œ ì‹¤íŒ¨: {del_err}")
                    return None
        
        # 3. "-en" íŒŒì¼ì´ ì—†ê±°ë‚˜ ì†ìƒë˜ì—ˆìœ¼ë©´ ìƒˆë¡œ ìƒì„±
        print(f"   ğŸ“ ì‘ì—… íŒŒì¼ì„ ìƒˆë¡œ ìƒì„±í•©ë‹ˆë‹¤...")
        return prepare_work_files(origin_path, completed_dir, completed_original, work_file_path)
        
    except Exception as e:
        print(f"   âŒ ì´ì–´í•˜ê¸° íŒŒì¼ ì¤€ë¹„ ì‹¤íŒ¨: {e}")
        return None


def main():
    """
    ë©”ì¸ í•¨ìˆ˜ - Google Sheets ê¸°ë°˜ ì‘ì—… ì²˜ë¦¬ì˜ ì§„ì…ì ì…ë‹ˆë‹¤.
    """
    print("=" * 60)
    print("ğŸŒ í•œí™”íì…€ ë²ˆì—­ í”„ë¡œì íŠ¸ v1.3.0")
    print("   Google Sheets ì—°ë™ ëŒ€ëŸ‰ ì²˜ë¦¬ ëª¨ë“œ")
    print("   .doc â†’ .docx ìë™ ë³€í™˜ ì§€ì›")
    print("   ğŸ“š ìš©ì–´ì§‘(Glossary) í”„ë¡¬í”„íŠ¸ ì£¼ì… ì§€ì›")
    print("=" * 60)
    
    # 1. ì„¤ì • ê²€ì¦
    is_valid, message = validate_config()
    if not is_valid:
        print(f"\nâŒ ì„¤ì • ì˜¤ë¥˜: {message}")
        return
    
    print("\nâœ… ì„¤ì • ê²€ì¦ ì™„ë£Œ")
    
    # 2. ìš©ì–´ì§‘ ë¡œë“œ
    glossary = get_glossary()
    if glossary.is_loaded:
        print(f"âœ… ìš©ì–´ì§‘ ë¡œë“œ ì™„ë£Œ: {glossary.get_term_count()}ê°œ ìš©ì–´")
    else:
        print("âš ï¸ ìš©ì–´ì§‘ ì—†ì´ ì§„í–‰í•©ë‹ˆë‹¤ (data/ìš©ì–´ì •ì˜.xlsx íŒŒì¼ í™•ì¸ í•„ìš”)")
    
    # 3. Google Sheets ì—°ê²°
    try:
        sheets_manager = SheetsManager()
    except Exception as e:
        print(f"\nâŒ Google Sheets ì—°ê²° ì‹¤íŒ¨: {e}")
        print("\nğŸ’¡ í•´ê²° ë°©ë²•:")
        print("   1. credentials.json íŒŒì¼ì´ í”„ë¡œì íŠ¸ ë£¨íŠ¸ì— ìˆëŠ”ì§€ í™•ì¸í•˜ì„¸ìš”.")
        print("   2. ì„œë¹„ìŠ¤ ê³„ì •ì— ìŠ¤í”„ë ˆë“œì‹œíŠ¸ ì ‘ê·¼ ê¶Œí•œì´ ìˆëŠ”ì§€ í™•ì¸í•˜ì„¸ìš”.")
        return
    
    # 4. ì‘ì—… ë£¨í”„ ì‹œì‘
    success_count = 0
    fail_count = 0
    
    # ì—°ì† ì˜¤ë¥˜ ë°©ì§€ë¥¼ ìœ„í•œ ì¶”ì  ë³€ìˆ˜ (íŒŒì¼ëª… ê¸°ë°˜)
    fail_count_by_file = {}  # {file_name: ì‹¤íŒ¨ íšŸìˆ˜}
    MAX_CONSECUTIVE_FAILS = 3  # ê°™ì€ íŒŒì¼ ì—°ì† ì‹¤íŒ¨ í—ˆìš© íšŸìˆ˜
    skipped_files = set()  # ê±´ë„ˆë›´ íŒŒì¼ëª… ëª©ë¡
    
    print("\nğŸš€ ì‘ì—… ì‹œì‘...")
    print("   (Ctrl+Cë¡œ ì¤‘ë‹¨í•  ìˆ˜ ìˆìŠµë‹ˆë‹¤)")
    
    while True:
        try:
            # ëŒ€ê¸° ì¤‘ì¸ ì‘ì—… ì¡°íšŒ
            task = sheets_manager.get_next_waiting_task()
            
            if task is None:
                print("\n" + "=" * 60)
                print("âœ… ëª¨ë“  ëŒ€ê¸° ì‘ì—… ì™„ë£Œ!")
                break
            
            current_row = task['row_index']
            file_name = task['file_name']
            
            # ì´ë¯¸ ê±´ë„ˆë›´ íŒŒì¼ì´ë©´ "ê±´ë„ˆë›°ê¸°ì™„ë£Œ"ë¡œ ë³€ê²½í•˜ê³  ë‹¤ìŒìœ¼ë¡œ
            if file_name in skipped_files:
                print(f"\nâ­ï¸ ê±´ë„ˆë›°ê¸°: {file_name} (ë°˜ë³µ ì‹¤íŒ¨ë¡œ ì œì™¸ë¨)")
                sheets_manager.update_status(current_row, "ê±´ë„ˆë›°ê¸°ì™„ë£Œ")
                continue
            
            # í•´ë‹¹ íŒŒì¼ì˜ ì‹¤íŒ¨ íšŸìˆ˜ í™•ì¸
            current_fail_count = fail_count_by_file.get(file_name, 0)
            
            if current_fail_count >= MAX_CONSECUTIVE_FAILS:
                print(f"\nâš ï¸ íŒŒì¼ '{file_name}'ì´(ê°€) {MAX_CONSECUTIVE_FAILS}íšŒ ì‹¤íŒ¨í–ˆìŠµë‹ˆë‹¤.")
                print(f"   â†’ ì´ íŒŒì¼ì„ ê±´ë„ˆë›°ê³  ë‹¤ìŒ íŒŒì¼ë¡œ ì´ë™í•©ë‹ˆë‹¤.")
                skipped_files.add(file_name)
                sheets_manager.update_status(current_row, "ê±´ë„ˆë›°ê¸°ì™„ë£Œ")
                fail_count += 1
                continue
            
            # ì‘ì—… ì²˜ë¦¬
            if process_task(sheets_manager, task):
                success_count += 1
                # ì„±ê³µ ì‹œ ì‹¤íŒ¨ ì¹´ìš´íŠ¸ ì œê±°
                if file_name in fail_count_by_file:
                    del fail_count_by_file[file_name]
            else:
                fail_count += 1
                # ì‹¤íŒ¨ ì‹œ ì¹´ìš´íŠ¸ ì¦ê°€
                fail_count_by_file[file_name] = current_fail_count + 1
                print(f"   âš ï¸ ì‹¤íŒ¨ íšŸìˆ˜: {fail_count_by_file[file_name]}/{MAX_CONSECUTIVE_FAILS}")
                
        except KeyboardInterrupt:
            print("\n\nâš ï¸ ì‚¬ìš©ìì— ì˜í•´ ì¤‘ë‹¨ë˜ì—ˆìŠµë‹ˆë‹¤.")
            break
        except Exception as e:
            print(f"\nâŒ ì˜ˆê¸°ì¹˜ ì•Šì€ ì˜¤ë¥˜: {e}")
            fail_count += 1
    
    # 5. ìµœì¢… ê²°ê³¼ ì¶œë ¥
    print("\n" + "=" * 60)
    print("ğŸ“Š ì‘ì—… ì™„ë£Œ ìš”ì•½")
    print("=" * 60)
    print(f"   âœ… ì„±ê³µ: {success_count}ê°œ")
    print(f"   âŒ ì‹¤íŒ¨: {fail_count}ê°œ")
    print(f"   ğŸ“ ê²°ê³¼ ìœ„ì¹˜: {COMPLETED_FOLDER}")
    print("=" * 60)


if __name__ == "__main__":
    main()
