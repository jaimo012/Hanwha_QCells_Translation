"""
PowerPoint 문서 핸들러 (PPTX Handler)

PowerPoint 문서(.pptx)의 번역을 처리합니다.
- 슬라이드, 그룹, 표 내부까지 재귀적 탐색
- 서식(폰트, 정렬, 색상 등)을 완벽하게 보존
- 토큰 사용량 추적
"""

import os
import time
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..config import BATCH_SIZE_PPTX, AUTO_SAVE_INTERVAL
from ..utils import has_korean
from ..translator import translate_batch


def copy_pptx_style_and_replace(paragraph, new_text):
    """
    텍스트를 교체하면서 문단의 모든 서식을 보존합니다.
    
    Args:
        paragraph: 대상 문단 객체
        new_text (str): 교체할 새 텍스트
    """
    style_backup = {}
    
    # 1. 문단 정렬(Alignment) 백업
    paragraph_alignment = paragraph.alignment 
    
    # 2. 폰트 스타일 백업 (첫 번째 Run 기준)
    if paragraph.runs:
        ref_run = paragraph.runs[0]
        style_backup = {
            'bold': ref_run.font.bold,
            'italic': ref_run.font.italic,
            'underline': ref_run.font.underline,
            'name': ref_run.font.name,
            'size': ref_run.font.size,
            'color': ref_run.font.color.rgb if hasattr(ref_run.font.color, 'rgb') else None
        }

    # 3. 텍스트 교체
    paragraph.clear() 
    new_run = paragraph.add_run()
    new_run.text = new_text
    
    # 4. 스타일 재적용
    if paragraph_alignment is not None:
        paragraph.alignment = paragraph_alignment
        
    if style_backup:
        new_run.font.bold = style_backup.get('bold')
        new_run.font.italic = style_backup.get('italic')
        new_run.font.underline = style_backup.get('underline')
        
        if style_backup.get('name'):
            new_run.font.name = style_backup['name']
        if style_backup.get('size'):
            new_run.font.size = style_backup['size']
        
        if style_backup.get('color'):
            try:
                new_run.font.color.rgb = style_backup['color']
            except:
                pass


def iter_pptx_shapes(shapes):
    """
    슬라이드의 모든 도형을 재귀적으로 순회하는 제너레이터입니다.
    
    Args:
        shapes: 슬라이드의 shapes 컬렉션
        
    Yields:
        shape 또는 cell: 텍스트를 포함하는 도형 또는 표 셀
    """
    for shape in shapes:
        # 1. 그룹(Group) -> 재귀 호출
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_pptx_shapes(shape.shapes)
        # 2. 표(Table) -> 셀 단위 탐색
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell
        # 3. 텍스트 프레임이 있는 일반 도형
        elif shape.has_text_frame:
            yield shape


def process_pptx(file_path, context, sheets_manager=None, row_index=None):
    """
    PowerPoint 문서를 번역합니다.
    
    Args:
        file_path (str): 원본 PPT 파일 경로
        context (str): 번역 지침 (Context)
        sheets_manager (SheetsManager, optional): 시트 관리자 (토큰 추적용)
        row_index (int, optional): 시트 행 번호
        
    Returns:
        str: 번역된 파일의 경로 (성공 시)
        None: 실패 시
    """
    print(f"📊 PPT 처리 중: {os.path.basename(file_path)}")
    
    prs = Presentation(file_path)
    # 작업 파일을 그대로 덮어쓰기 (main.py에서 이미 " - en" 파일 생성)
    new_path = file_path
    
    batch_queue = []
    total_count = 0
    batch_cycle = 0
    total_input_tokens = 0
    total_output_tokens = 0

    for slide_idx, slide in enumerate(prs.slides):
        for shape_or_cell in iter_pptx_shapes(slide.shapes):
            
            text_frame = None
            
            if hasattr(shape_or_cell, 'has_text_frame'):
                if shape_or_cell.has_text_frame:
                    text_frame = shape_or_cell.text_frame
            elif hasattr(shape_or_cell, 'text_frame'):
                text_frame = shape_or_cell.text_frame
            
            if not text_frame:
                continue

            for p in text_frame.paragraphs:
                # p.text가 None일 수 있으므로 안전하게 처리
                text = p.text
                if text is None:
                    continue
                text = text.strip()
                
                if text and has_korean(text):
                    batch_queue.append(p)
                    
                    if len(batch_queue) >= BATCH_SIZE_PPTX:
                        texts = [b.text for b in batch_queue]
                        translated, input_tokens, output_tokens = translate_batch(texts, context)
                        
                        total_input_tokens += input_tokens
                        total_output_tokens += output_tokens
                        
                        if len(translated) == len(batch_queue):
                            for obj, trans_text in zip(batch_queue, translated):
                                copy_pptx_style_and_replace(obj, trans_text)
                            total_count += len(translated)
                            batch_cycle += 1
                            
                            # 실시간 진행 상황 (배치 횟수만 표시)
                            print(f"   ▶ 배치 {batch_cycle}회 진행 중...          ", end="\r")
                            
                            # 중간 저장
                            if batch_cycle % AUTO_SAVE_INTERVAL == 0:
                                print()  # 줄바꿈
                                print(f"   💾 [자동저장] 중간 저장 중...")
                                prs.save(new_path)
                                
                                # 시트에 토큰 사용량 업데이트
                                if sheets_manager and row_index:
                                    sheets_manager.update_tokens(
                                        row_index,
                                        total_input_tokens,
                                        total_output_tokens
                                    )
                                    total_input_tokens = 0
                                    total_output_tokens = 0
                        
                        batch_queue = []
                        time.sleep(0.5)

    # 잔여 배치 처리 (배치 크기보다 적은 남은 데이터)
    if batch_queue:
        print(f"\n   🔄 잔여 {len(batch_queue)}개 처리 중...")
        texts = [b.text for b in batch_queue]
        translated, input_tokens, output_tokens = translate_batch(texts, context)
        
        total_input_tokens += input_tokens
        total_output_tokens += output_tokens
        
        if len(translated) == len(batch_queue):
            for obj, trans_text in zip(batch_queue, translated):
                copy_pptx_style_and_replace(obj, trans_text)
            total_count += len(translated)
            batch_cycle += 1
        print(f"   ✅ 잔여 처리 완료")

    print(f"\n   💾 최종 저장 중...")
    prs.save(new_path)
    print(f"   ✅ 파일 저장 완료")
    
    # 최종 토큰 사용량 업데이트
    if sheets_manager and row_index:
        if total_input_tokens > 0 or total_output_tokens > 0:
            sheets_manager.update_tokens(
                row_index,
                total_input_tokens,
                total_output_tokens
            )
    
    print()  # 진행 상황 줄 종료
    print(f"   ✅ PPT 번역 완료: {batch_cycle}개 배치, {total_count}개 문단")
    return new_path
