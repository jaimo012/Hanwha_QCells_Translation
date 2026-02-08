"""
검수 모듈 (Verification Module)

1차 번역이 완료된 파일들을 검수하여 남아있는 한글을 추가 번역합니다.
- Google Sheets에서 "완료" 상태인 파일을 조회
- completed_folder의 "-en" 파일을 열어서 번역 프로세스 재실행
- 남아있는 한글이 있다면 번역하여 저장
- 진행상태를 "1차 검수완료"로 변경
"""

import os
import traceback

from docx import Document
from docx.text.paragraph import Paragraph
from pptx import Presentation
import openpyxl

from .config import (
    COMPLETED_FOLDER,
    SUPPORTED_EXTENSIONS,
    validate_config
)
from .translator import generate_context
from .handlers import process_docx, process_pptx, process_xlsx
from .sheets_manager import SheetsManager, Status
from .slack_notifier import send_review_completion_notification, send_error_notification
from .glossary import get_glossary
from .utils import has_korean


def build_work_file_path(upper_path, sub_path, file_name):
    """
    완료된 파일의 "-en" 작업 파일 경로를 구성합니다.
    
    Args:
        upper_path (str): 상위 경로 (예: "MC")
        sub_path (str): 세부 경로 (예: "10.분석단계")
        file_name (str): 원본 파일명 (확장자 포함)
        
    Returns:
        str: 작업 파일 경로 ("-en"이 붙은 파일)
    """
    # 확장자 추출
    name, ext = os.path.splitext(file_name)
    ext_lower = ext.lower()
    
    # .doc 파일은 .docx로 변환되어 작업됨
    if ext_lower == '.doc':
        work_file_name = f"{name} - en.docx"
    else:
        work_file_name = f"{name} - en{ext_lower}"
    
    # completed_folder 내 경로 구성
    work_file_path = os.path.join(COMPLETED_FOLDER, upper_path, sub_path, work_file_name)
    
    return work_file_path


def verify_file_exists(file_path):
    """
    파일이 존재하는지 확인합니다.
    
    Args:
        file_path (str): 확인할 파일 경로
        
    Returns:
        bool: 파일이 존재하면 True
    """
    return os.path.exists(file_path)


def verify_file_integrity(file_path):
    """
    파일의 무결성을 검사합니다.
    
    Args:
        file_path (str): 검사할 파일 경로
        
    Returns:
        bool: 파일이 정상이면 True
    """
    try:
        if os.path.getsize(file_path) == 0:
            return False
        
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.docx':
            doc = Document(file_path)
            _ = doc.paragraphs
            return True
        elif ext == '.pptx':
            prs = Presentation(file_path)
            _ = prs.slides
            return True
        elif ext == '.xlsx':
            return os.path.getsize(file_path) > 1000
        
        return True
        
    except Exception as e:
        print(f"   ⚠️ 파일 무결성 검사 실패: {e}")
        return False


def extract_sample_text(file_path):
    """
    파일에서 Context 분석용 샘플 텍스트를 추출합니다.
    
    Args:
        file_path (str): 파일 경로
        
    Returns:
        str: 추출된 샘플 텍스트
    """
    sample_text = ""
    file_path_lower = file_path.lower()
    
    try:
        if file_path_lower.endswith('.docx'):
            doc = Document(file_path)
            sample_text = "\n".join([p.text for p in doc.paragraphs[:300]])
            
        elif file_path_lower.endswith('.pptx'):
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                if i >= 3:
                    break
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        sample_text += shape.text + "\n"
                        
        elif file_path_lower.endswith('.xlsx'):
            sample_text = "MES Excel Data"
    except Exception as e:
        print(f"   ⚠️ 샘플 텍스트 추출 실패: {e}")
        sample_text = "MES Technical Document"
    
    return sample_text


# ==============================================================================
# [한글 스캔 함수들] - 번역 필요 여부를 미리 확인
# ==============================================================================

def scan_korean_in_docx(file_path):
    """
    Word 문서에서 한글이 있는지 스캔합니다.
    
    Args:
        file_path (str): 파일 경로
        
    Returns:
        tuple: (한글 존재 여부, 한글이 포함된 텍스트 개수)
    """
    korean_count = 0
    doc = None
    
    try:
        doc = Document(file_path)
        
        # 1. 본문 문단 스캔
        for p in doc.paragraphs:
            text = p.text
            if text and has_korean(text.strip()):
                korean_count += 1
        
        # 2. 표 내부 스캔
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        text = p.text
                        if text and has_korean(text.strip()):
                            korean_count += 1
        
        # 3. 텍스트 상자 스캔
        for element in doc.element.body.iter():
            if element.tag.endswith('txbxContent'):
                for child in element.iter():
                    if child.tag.endswith('p'):
                        p = Paragraph(child, doc)
                        text = p.text
                        if text and has_korean(text.strip()):
                            korean_count += 1
        
    except Exception as e:
        print(f"   ⚠️ Word 스캔 오류: {e}")
        return False, 0
    finally:
        # 메모리에서 명시적 해제 (파일 객체 정리)
        if doc is not None:
            del doc
    
    return korean_count > 0, korean_count


def scan_korean_in_pptx(file_path):
    """
    PowerPoint 문서에서 한글이 있는지 스캔합니다.
    
    Args:
        file_path (str): 파일 경로
        
    Returns:
        tuple: (한글 존재 여부, 한글이 포함된 텍스트 개수)
    """
    korean_count = 0
    prs = None
    
    try:
        prs = Presentation(file_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                # 일반 텍스트 shape
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text
                            if text and has_korean(text.strip()):
                                korean_count += 1
                
                # 표(Table) 내부
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if hasattr(cell, "text_frame"):
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        text = run.text
                                        if text and has_korean(text.strip()):
                                            korean_count += 1
                
                # 그룹 shape 내부
                if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                    for sub_shape in shape.shapes:
                        if hasattr(sub_shape, "text_frame"):
                            for paragraph in sub_shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text = run.text
                                    if text and has_korean(text.strip()):
                                        korean_count += 1
        
    except Exception as e:
        print(f"   ⚠️ PowerPoint 스캔 오류: {e}")
        return False, 0
    finally:
        # 메모리에서 명시적 해제 (파일 객체 정리)
        if prs is not None:
            del prs
    
    return korean_count > 0, korean_count


def scan_korean_in_xlsx(file_path):
    """
    Excel 문서에서 한글이 있는지 스캔합니다.
    
    Args:
        file_path (str): 파일 경로
        
    Returns:
        tuple: (한글 존재 여부, 한글이 포함된 셀 개수)
    """
    korean_count = 0
    wb = None
    
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        if has_korean(cell.value.strip()):
                            korean_count += 1
        
    except Exception as e:
        print(f"   ⚠️ Excel 스캔 오류: {e}")
        return False, 0
    finally:
        # 성공/실패와 무관하게 반드시 워크북 닫기
        if wb is not None:
            wb.close()
    
    return korean_count > 0, korean_count


def scan_korean_in_file(file_path):
    """
    파일 형식에 따라 한글 존재 여부를 스캔합니다.
    
    Args:
        file_path (str): 파일 경로
        
    Returns:
        tuple: (한글 존재 여부, 한글이 포함된 항목 개수)
    """
    file_path_lower = file_path.lower()
    
    if file_path_lower.endswith('.docx'):
        return scan_korean_in_docx(file_path)
    elif file_path_lower.endswith('.pptx'):
        return scan_korean_in_pptx(file_path)
    elif file_path_lower.endswith('.xlsx'):
        return scan_korean_in_xlsx(file_path)
    else:
        print(f"   ⚠️ 지원하지 않는 파일 형식: {file_path}")
        return False, 0


def process_single_file(work_file_path, file_context, sheets_manager=None, row_index=None):
    """
    단일 파일을 번역 처리합니다.
    
    Args:
        work_file_path (str): 작업 파일 경로
        file_context (str): 번역 지침 (Context)
        sheets_manager (SheetsManager, optional): 시트 관리자
        row_index (int, optional): 시트 행 번호
        
    Returns:
        str: 번역된 파일 경로 (성공 시)
        None: 실패 시
    """
    file_path_lower = work_file_path.lower()
    
    if file_path_lower.endswith('.docx'):
        return process_docx(work_file_path, file_context, sheets_manager, row_index)
    elif file_path_lower.endswith('.pptx'):
        return process_pptx(work_file_path, file_context, sheets_manager, row_index)
    elif file_path_lower.endswith('.xlsx'):
        return process_xlsx(work_file_path, file_context, sheets_manager, row_index)
    else:
        print(f"   ⚠️ 지원하지 않는 파일 형식: {work_file_path}")
        return None


def verify_task(sheets_manager, task):
    """
    단일 검수 작업을 처리합니다.
    
    [최적화] 한글이 있는지 먼저 스캔 후, 있을 때만 Context 분석 → 번역 진행
    
    - completed_folder에서 "-en" 파일을 찾아서
    - 한글이 있는지 먼저 스캔 (API 호출 없음)
    - 한글이 있으면: Context 분석 → 번역 프로세스 실행
    - 한글이 없으면: 바로 "1차 검수완료" 처리 (빠른 처리)
    - 진행상태를 "1차 검수완료"로 변경
    
    Args:
        sheets_manager (SheetsManager): 시트 관리자
        task (dict): 작업 정보
        
    Returns:
        tuple: (성공 여부, 결과 타입)
            - (True, "translated"): 번역 수행 후 완료
            - (True, "no_korean"): 한글 없음, 바로 완료
            - (False, "error"): 오류 발생
    """
    row_index = task['row_index']
    upper_path = task['upper_path']
    sub_path = task['sub_path']
    file_name = task['file_name']
    
    print(f"\n{'='*60}")
    print(f"🔍 검수 시작: {file_name}")
    print(f"   경로: {upper_path}/{sub_path}")
    print(f"{'='*60}")
    
    try:
        # 1. "-en" 작업 파일 경로 구성
        work_file_path = build_work_file_path(upper_path, sub_path, file_name)
        
        # 2. 파일 존재 확인
        if not verify_file_exists(work_file_path):
            raise FileNotFoundError(f"작업 파일을 찾을 수 없습니다: {work_file_path}")
        
        print(f"   ✅ 작업 파일 발견: {os.path.basename(work_file_path)}")
        
        # 3. 파일 무결성 검사
        if not verify_file_integrity(work_file_path):
            raise ValueError(f"작업 파일이 손상되었습니다: {work_file_path}")
        
        print(f"   ✅ 파일 무결성 확인 완료")
        
        # 4. [최적화] 한글 존재 여부 먼저 스캔 (API 호출 없음)
        print(f"   🔎 한글 스캔 중...")
        has_korean_text, korean_count = scan_korean_in_file(work_file_path)
        
        if not has_korean_text:
            # ✨ 한글이 없으면 바로 검수 완료 처리 (Context 분석 스킵)
            print(f"   ✨ 한글 없음 → 번역 불필요, 바로 검수 완료 처리")
            sheets_manager.update_status(row_index, Status.REVIEW_1_COMPLETED)
            print(f"\n   🎉 검수 완료! (번역 대상 없음)")
            return True, "no_korean"
        
        print(f"   📝 한글 발견: {korean_count}개 항목 → 번역 필요")
        
        # 5. 상태를 "진행중"으로 변경 (검수 중임을 표시)
        sheets_manager.update_status(row_index, Status.IN_PROGRESS)
        print(f"   🔄 상태 변경: 진행중 (검수)")
        
        # 6. Context 분석 (한글이 있을 때만 실행)
        print(f"   🤖 Context 분석 중...")
        sample_text = extract_sample_text(work_file_path)
        file_context = generate_context(sample_text)
        print(f"   ✅ Context 분석 완료")
        
        # 7. 번역 프로세스 실행 (남아있는 한글 번역)
        # 핸들러는 has_korean()으로 한글이 있는 텍스트만 번역하므로
        # 이미 번역된 부분은 건너뛰고 남은 한글만 번역함
        result = process_single_file(work_file_path, file_context, sheets_manager, row_index)
        
        if result:
            # 8. 검수 완료 처리
            sheets_manager.update_status(row_index, Status.REVIEW_1_COMPLETED)
            print(f"\n   🎉 검수 완료!")
            
            # 9. Slack 알림 전송 (검수 완료 전용 포맷)
            try:
                review_progress = sheets_manager.get_review_progress()
                file_path = f"{upper_path}/{sub_path}"
                
                # "-en"이 붙은 실제 작업 파일명 사용
                work_file_name = os.path.basename(work_file_path)
                
                send_review_completion_notification(
                    file_name=work_file_name,
                    file_path=file_path,
                    review_progress_percent=review_progress
                )
            except Exception as slack_error:
                print(f"   ⚠️ Slack 알림 전송 실패: {slack_error}")
            
            return True, "translated"
        else:
            raise Exception("검수 처리 실패")
            
    except Exception as e:
        # 오류 발생 시 - 상태를 "완료"로 되돌리고 오류 기록
        error_msg = str(e)
        error_trace = traceback.format_exc()
        module_name = "verify.verify_task"
        
        detailed_error = f"[검수오류] {error_msg}\n\n상세:\n{error_trace}"
        
        print(f"\n   ❌ 검수 오류 발생: {error_msg}")
        
        # 상태를 원래대로 "완료"로 되돌림 (다음 검수 시 재시도 가능)
        sheets_manager.update_status(row_index, Status.COMPLETED)
        
        # 오류 내용은 비고(K열)에 기록 (J열은 원본 번역 오류용)
        try:
            sheets_manager.sheet.update_cell(row_index, 11, detailed_error[:500])  # K열 = 11
        except Exception:
            pass
        
        # Slack 오류 알림
        try:
            slack_error_msg = f"*파일*: {file_name}\n*경로*: {upper_path}/{sub_path}\n*검수오류*: {error_msg}"
            send_error_notification(slack_error_msg)
        except Exception:
            pass
        
        return False, "error"


def main():
    """
    검수 프로세스 메인 함수
    
    Google Sheets에서 "완료" 상태인 파일들을 조회하여
    순차적으로 검수를 진행합니다.
    """
    print("=" * 60)
    print("🔍 한화큐셀 번역 프로젝트 - 1차 검수 프로세스")
    print("   완료된 파일의 남은 한글을 추가 번역합니다")
    print("=" * 60)
    
    # 1. 설정 검증
    is_valid, message = validate_config()
    if not is_valid:
        print(f"\n❌ 설정 오류: {message}")
        return
    
    print("\n✅ 설정 검증 완료")
    
    # 2. 용어집 로드
    glossary = get_glossary()
    if glossary.is_loaded:
        print(f"✅ 용어집 로드 완료: {glossary.get_term_count()}개 용어")
    else:
        print("⚠️ 용어집 없이 진행합니다")
    
    # 3. Google Sheets 연결
    try:
        sheets_manager = SheetsManager()
    except Exception as e:
        print(f"\n❌ Google Sheets 연결 실패: {e}")
        return
    
    # 4. "완료" 상태인 작업 목록 조회
    completed_tasks = sheets_manager.get_completed_tasks()
    
    if not completed_tasks:
        print("\n" + "=" * 60)
        print("✅ 검수할 파일이 없습니다. (완료 상태 파일 0개)")
        print("=" * 60)
        return
    
    print(f"\n📋 검수 대상: {len(completed_tasks)}개 파일")
    
    # 5. 검수 루프 시작
    success_count = 0
    fail_count = 0
    skip_count = 0        # 반복 실패로 건너뛴 파일
    no_korean_count = 0   # 번역 불필요 (한글 없음)
    
    # 연속 오류 방지
    fail_count_by_file = {}
    MAX_CONSECUTIVE_FAILS = 3
    skipped_files = set()
    
    print("\n🚀 검수 시작...")
    print("   (Ctrl+C로 중단할 수 있습니다)")
    
    for task in completed_tasks:
        try:
            file_name = task['file_name']
            
            # 이미 건너뛴 파일이면 스킵
            if file_name in skipped_files:
                print(f"\n⏭️ 건너뛰기: {file_name} (반복 실패로 제외됨)")
                skip_count += 1
                continue
            
            # 실패 횟수 확인
            current_fail_count = fail_count_by_file.get(file_name, 0)
            
            if current_fail_count >= MAX_CONSECUTIVE_FAILS:
                print(f"\n⚠️ 파일 '{file_name}'이(가) {MAX_CONSECUTIVE_FAILS}회 실패했습니다.")
                print(f"   → 이 파일을 건너뛰고 다음 파일로 이동합니다.")
                skipped_files.add(file_name)
                skip_count += 1
                continue
            
            # 검수 실행
            success, result_type = verify_task(sheets_manager, task)
            
            if success:
                if result_type == "no_korean":
                    no_korean_count += 1  # 한글 없음 (번역 불필요)
                else:
                    success_count += 1    # 번역 수행 완료
                
                # 성공 시 실패 카운트 제거
                if file_name in fail_count_by_file:
                    del fail_count_by_file[file_name]
            else:
                fail_count += 1
                fail_count_by_file[file_name] = current_fail_count + 1
                print(f"   ⚠️ 실패 횟수: {fail_count_by_file[file_name]}/{MAX_CONSECUTIVE_FAILS}")
                
        except KeyboardInterrupt:
            print("\n\n⚠️ 사용자에 의해 중단되었습니다.")
            break
        except Exception as e:
            print(f"\n❌ 예기치 않은 오류: {e}")
            fail_count += 1
    
    # 6. 최종 결과 출력
    print("\n" + "=" * 60)
    print("📊 검수 완료 요약")
    print("=" * 60)
    print(f"   ✅ 번역 완료: {success_count}개 (한글 → 영어 추가 번역)")
    print(f"   ✨ 번역 불필요: {no_korean_count}개 (한글 없음)")
    print(f"   ❌ 검수 실패: {fail_count}개")
    print(f"   ⏭️ 건너뜀: {skip_count}개")
    print(f"   📁 파일 위치: {COMPLETED_FOLDER}")
    print("=" * 60)
    
    # API 비용 절감 효과 표시
    if no_korean_count > 0:
        print(f"\n💡 최적화 효과: {no_korean_count}개 파일에서 Context 분석 API 호출 생략")


if __name__ == "__main__":
    main()
